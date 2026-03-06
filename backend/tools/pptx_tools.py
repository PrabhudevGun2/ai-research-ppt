import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from typing import List, Optional
import logging

logger = logging.getLogger(__name__)

# Color scheme
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)   # Deep indigo header
ACCENT_BLUE = RGBColor(0x19, 0x76, 0xD2)  # Medium blue accents
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
BULLET_COLOR = RGBColor(0x19, 0x76, 0xD2)

SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation() -> Presentation:
    """Create a new presentation with widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def set_slide_background(slide, color: RGBColor = LIGHT_GRAY):
    """Set a solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: Optional[str] = None):
    """Add a dark header bar with title text."""
    # Header rectangle
    header = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(1.4),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Title text box
    txBox = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.15),
        Inches(12.0), Inches(0.9),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.4), Inches(1.05),
            Inches(12.0), Inches(0.35),
        )
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = ACCENT_BLUE
        sp.font.italic = True


def add_bullet_points(slide, points: List[str], top: float = 1.6, left: float = 0.5,
                       width: float = 12.3, height: float = 5.4, font_size: int = 18):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)


def add_two_column_layout(slide, left_points: List[str], right_points: List[str],
                           top: float = 1.6):
    """Add two columns of bullet points."""
    # Left column
    add_bullet_points(slide, left_points, top=top, left=0.4, width=6.0, font_size=16)
    # Right column
    add_bullet_points(slide, right_points, top=top, left=6.8, width=6.0, font_size=16)


def add_accent_line(slide, top: float = 1.35):
    """Add a thin accent line below the header."""
    line = slide.shapes.add_shape(
        1,
        Inches(0), Inches(top),
        SLIDE_WIDTH, Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()


def add_page_number(slide, number: int, total: int):
    """Add a subtle page number in the bottom-right corner."""
    txBox = slide.shapes.add_textbox(
        Inches(11.8), Inches(7.1),
        Inches(1.3), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{number} / {total}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x9E, 0x9E, 0x9E)


def add_footer(slide, text: str):
    """Add a footer bar at the bottom of the slide."""
    footer = slide.shapes.add_shape(
        1,
        Inches(0), Inches(7.2),
        SLIDE_WIDTH, Inches(0.3),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = DARK_BLUE
    footer.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.2),
        Inches(10.0), Inches(0.3),
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE


def add_image_to_slide(slide, image_path: str, left: float = Inches(1), top: float = Inches(2),
                       max_width: float = Inches(10), max_height: float = Inches(5)):
    """Add an image to a slide with automatic scaling to fit constraints."""
    from PIL import Image as PILImage

    if not os.path.exists(image_path):
        logger.warning(f"Image not found: {image_path}")
        return

    # Get original image dimensions
    with PILImage.open(image_path) as img:
        orig_width, orig_height = img.size

    # Calculate scaling factor to fit within max dimensions
    width_scale = max_width / Inches(1) / (orig_width / 96)  # 96 DPI assumption
    height_scale = max_height / Inches(1) / (orig_height / 96)
    scale = min(width_scale, height_scale, 1.0)  # Don't upscale

    # Calculate final dimensions
    final_width = int(orig_width * scale)
    final_height = int(orig_height * scale)

    # Convert to EMUs (English Metric Units)
    width_emu = int(final_width * 914400 / 96)  # 914400 EMUs per inch at 96 DPI
    height_emu = int(final_height * 914400 / 96)

    try:
        slide.shapes.add_picture(image_path, left, top, width_emu, height_emu)
        logger.info(f"Added image to slide: {image_path}")
    except Exception as e:
        logger.error(f"Failed to add image to slide: {e}")
