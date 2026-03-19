"""
PPT Generation agent - creates professional PowerPoint with images and tables.
Also generates a companion Word document with full verbose content.
"""
import os
import logging
from datetime import date

from backend.graph.state import ResearchState, Stage, SlideContent, GeneratedPPT
from backend.tools.pptx_tools import (
    create_presentation, set_slide_background, add_header_bar,
    add_bullet_points, add_accent_line, add_page_number, add_footer,
    add_image_to_slide, DARK_BLUE, WHITE, LIGHT_GRAY, ACCENT_BLUE,
)
from backend.config import settings
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# Color scheme
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PROBLEM_RED = RGBColor(0xC6, 0x28, 0x28)
METHOD_GREEN = RGBColor(0x2E, 0x7D, 0x32)
RESULTS_BLUE = RGBColor(0x15, 0x65, 0xC0)


def _add_title_slide(prs, content: SlideContent, num: int, total: int):
    """Title slide - clean and professional."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.33), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content["title"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if content.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.33), Inches(1.0))
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = content["subtitle"]
        sp.font.size = Pt(20)
        sp.font.color.rgb = ACCENT_BLUE
        sp.alignment = PP_ALIGN.CENTER

    # Date and arxiv info
    if content.get("body_points"):
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.33), Inches(1.0))
        itf = info_box.text_frame
        for i, point in enumerate(content["body_points"][:2]):
            ip = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
            ip.text = point
            ip.font.size = Pt(14)
            ip.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)
            ip.alignment = PP_ALIGN.CENTER

    add_page_number(slide, num, total)


def _add_content_slide_with_image(prs, content: SlideContent, num: int, total: int, accent_color=None):
    """Content slide with optional image on the right side."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, LIGHT_GRAY)

    add_header_bar(slide, content["title"], content.get("subtitle"))

    # Add accent line
    if accent_color:
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.45), Inches(12.33), Pt(4))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_color
        line.line.fill.background()

    # Check if we have an image
    image_path = content.get("image_path")
    has_valid_image = image_path and os.path.exists(image_path)

    if has_valid_image:
        # Two-column layout: text on left, image on right
        # NOTE: add_bullet_points expects raw floats (inches), NOT Inches() objects
        add_bullet_points(
            slide, content["body_points"],
            font_size=14, top=1.7,
            left=0.4, width=6.5
        )

        # add_image_to_slide expects Inches() objects for positioning
        try:
            add_image_to_slide(slide, image_path, left=Inches(7.2), top=Inches(1.7),
                             max_width=Inches(5.5), max_height=Inches(5.0))
        except Exception as e:
            logger.warning(f"Failed to add image {image_path}: {e}")
    else:
        # Full-width text
        add_bullet_points(slide, content["body_points"], font_size=16, top=1.8)

    add_page_number(slide, num, total)

    footer_text = content.get("topic") or content.get("subtitle") or "Technical Presentation"
    if footer_text and isinstance(footer_text, str):
        add_footer(slide, footer_text[:50])


def _add_problem_slide(prs, content: SlideContent, num: int, total: int):
    _add_content_slide_with_image(prs, content, num, total, accent_color=PROBLEM_RED)


def _add_methodology_slide(prs, content: SlideContent, num: int, total: int):
    _add_content_slide_with_image(prs, content, num, total, accent_color=METHOD_GREEN)


def _add_results_slide(prs, content: SlideContent, num: int, total: int):
    _add_content_slide_with_image(prs, content, num, total, accent_color=RESULTS_BLUE)


def _add_equation_slide(prs, content: SlideContent, num: int, total: int):
    """Equation slide - prominent, centered display."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, LIGHT_GRAY)

    add_header_bar(slide, content["title"], "Mathematical Formulation")

    image_path = content.get("image_path")
    if image_path and os.path.exists(image_path):
        try:
            add_image_to_slide(slide, image_path, left=Inches(1.0), top=Inches(2.0),
                             max_width=Inches(11.0), max_height=Inches(3.5))
            if content["body_points"]:
                add_bullet_points(slide, content["body_points"], font_size=14,
                                top=5.5, left=0.5, width=12.0)
        except Exception as e:
            logger.warning(f"Failed to add equation image: {e}")
            _add_equation_text(slide, content)
    else:
        _add_equation_text(slide, content)

    add_page_number(slide, num, total)
    add_footer(slide, "Key Equations")


def _add_equation_text(slide, content: SlideContent):
    """Display equations as text (fallback)."""
    y_pos = Inches(1.8)
    for point in content["body_points"][:5]:
        eq_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.33), Inches(0.8))
        tf = eq_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        clean_text = point.replace("$", "").replace("\\(", "").replace("\\)", "")
        p.text = clean_text
        p.font.size = Pt(18)
        p.font.name = "Courier New"
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        y_pos += Inches(0.7)


def _add_architecture_slide(prs, content: SlideContent, num: int, total: int):
    """Architecture slide - often has diagram."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, LIGHT_GRAY)

    add_header_bar(slide, content["title"], "System Architecture")

    image_path = content.get("image_path")
    if image_path and os.path.exists(image_path):
        try:
            add_image_to_slide(slide, image_path, left=Inches(1.0), top=Inches(1.6),
                             max_width=Inches(11.0), max_height=Inches(4.5))
            if content["body_points"]:
                add_bullet_points(slide, content["body_points"][:3], font_size=13,
                                top=6.2, left=0.5, width=12.0)
        except Exception as e:
            logger.warning(f"Failed to add architecture image: {e}")
            _add_content_slide_with_image(prs, content, num, total)
    else:
        _add_content_slide_with_image(prs, content, num, total)

    add_page_number(slide, num, total)
    add_footer(slide, "Architecture")


def _add_summary_slide(prs, content: SlideContent, num: int, total: int):
    """Summary/conclusion slide with dark background."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, DARK_BLUE)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = content["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if content["body_points"]:
        points_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(4.0))
        ptf = points_box.text_frame
        for i, point in enumerate(content["body_points"][:8]):
            pp = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
            pp.text = f"• {point}"
            pp.font.size = Pt(16)
            pp.font.color.rgb = ACCENT_BLUE
            pp.space_after = Pt(10)

    add_page_number(slide, num, total)


# Slide type handlers
SLIDE_HANDLERS = {
    "title": _add_title_slide,
    "problem": _add_problem_slide,
    "background": _add_content_slide_with_image,
    "contribution": _add_content_slide_with_image,
    "methodology": _add_methodology_slide,
    "method": _add_methodology_slide,
    "architecture": _add_architecture_slide,
    "equation": _add_equation_slide,
    "equations": _add_equation_slide,
    "algorithm": _add_content_slide_with_image,
    "experiments": _add_content_slide_with_image,
    "results": _add_results_slide,
    "evaluation": _add_results_slide,
    "analysis": _add_content_slide_with_image,
    "discussion": _add_content_slide_with_image,
    "limitations": _add_content_slide_with_image,
    "future": _add_content_slide_with_image,
    "conclusion": _add_summary_slide,
    "references": _add_content_slide_with_image,
    "content": _add_content_slide_with_image,
}


def _generate_word_document(session_id: str, slides: list, paper: dict) -> str:
    """Generate a professional Word document companion to the PPT."""
    from docx import Document
    from docx.shared import Pt as DocxPt, Inches as DocxInches, RGBColor as DocxRGB, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn

    doc = Document()

    # -- Page setup: narrower margins for professional look --------------------
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    # -- Restyle built-in styles -----------------------------------------------
    style_normal = doc.styles['Normal']
    style_normal.font.size = DocxPt(11)
    style_normal.font.name = 'Calibri'
    style_normal.paragraph_format.space_after = DocxPt(6)
    style_normal.paragraph_format.line_spacing = 1.15

    # Heading 1 – dark blue, bold, 18pt
    h1 = doc.styles['Heading 1']
    h1.font.size = DocxPt(18)
    h1.font.bold = True
    h1.font.color.rgb = DocxRGB(0x1A, 0x23, 0x7E)
    h1.font.name = 'Calibri'
    h1.paragraph_format.space_before = DocxPt(24)
    h1.paragraph_format.space_after = DocxPt(8)

    # Heading 2 – accent blue, 14pt
    h2 = doc.styles['Heading 2']
    h2.font.size = DocxPt(14)
    h2.font.bold = True
    h2.font.color.rgb = DocxRGB(0x19, 0x76, 0xD2)
    h2.font.name = 'Calibri'
    h2.paragraph_format.space_before = DocxPt(16)
    h2.paragraph_format.space_after = DocxPt(6)

    title_text = paper.get("title", "Research Paper Analysis")
    authors = paper.get("authors", [])
    arxiv_id = paper.get("arxiv_id", "")
    abstract = paper.get("abstract", "")

    # ====== TITLE PAGE ========================================================
    # Add vertical spacing to push title down
    for _ in range(6):
        doc.add_paragraph()

    # Thin accent line above title
    line_para = doc.add_paragraph()
    line_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = line_para.add_run("━" * 50)
    run.font.color.rgb = DocxRGB(0x19, 0x76, 0xD2)
    run.font.size = DocxPt(10)

    # Title
    tp = doc.add_paragraph()
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = tp.add_run(title_text)
    tr.font.size = DocxPt(28)
    tr.bold = True
    tr.font.color.rgb = DocxRGB(0x1A, 0x23, 0x7E)
    tr.font.name = 'Calibri'

    # Subtitle: "Research Paper Analysis"
    sp = doc.add_paragraph()
    sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = sp.add_run("Research Paper Analysis")
    sr.font.size = DocxPt(14)
    sr.font.color.rgb = DocxRGB(0x75, 0x75, 0x75)
    sr.font.italic = True

    # Thin accent line below subtitle
    line_para2 = doc.add_paragraph()
    line_para2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = line_para2.add_run("━" * 50)
    run2.font.color.rgb = DocxRGB(0x19, 0x76, 0xD2)
    run2.font.size = DocxPt(10)

    doc.add_paragraph()  # spacing

    # Authors
    if authors:
        ap = doc.add_paragraph()
        ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ar = ap.add_run(', '.join(authors[:10]))
        ar.font.size = DocxPt(12)
        ar.font.color.rgb = DocxRGB(0x44, 0x44, 0x44)

    # Metadata line
    meta_parts = []
    if arxiv_id:
        meta_parts.append(f"ArXiv: {arxiv_id}")
    meta_parts.append(f"Generated: {date.today().strftime('%B %d, %Y')}")
    mp = doc.add_paragraph()
    mp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    mr = mp.add_run(" | ".join(meta_parts))
    mr.font.size = DocxPt(10)
    mr.font.color.rgb = DocxRGB(0x99, 0x99, 0x99)

    doc.add_page_break()

    # ====== TABLE OF CONTENTS =================================================
    doc.add_heading("Table of Contents", level=1)
    section_num = 0
    for slide in slides:
        if slide.get("slide_type") == "title":
            continue
        section_num += 1
        toc_para = doc.add_paragraph()
        toc_para.paragraph_format.space_after = DocxPt(3)
        num_run = toc_para.add_run(f"{section_num}. ")
        num_run.font.bold = True
        num_run.font.color.rgb = DocxRGB(0x19, 0x76, 0xD2)
        title_run = toc_para.add_run(slide.get('title', f'Section {section_num}'))
        title_run.font.size = DocxPt(11)

    doc.add_page_break()

    # ====== ABSTRACT ==========================================================
    if abstract:
        doc.add_heading("Abstract", level=1)
        abs_para = doc.add_paragraph()
        abs_para.paragraph_format.first_line_indent = Cm(1)
        abs_para.paragraph_format.space_after = DocxPt(12)
        abs_run = abs_para.add_run(abstract)
        abs_run.font.size = DocxPt(11)
        abs_run.font.color.rgb = DocxRGB(0x33, 0x33, 0x33)
        doc.add_paragraph()

    # ====== SECTIONS FROM SLIDES ==============================================
    section_num = 0
    for slide in slides:
        if slide.get("slide_type") == "title":
            continue

        section_num += 1
        slide_title = slide.get("title", "Section")
        slide_type = slide.get("slide_type", "content")

        # Section heading
        doc.add_heading(f"{section_num}. {slide_title}", level=1)

        # Subtitle as a styled intro
        subtitle = slide.get("subtitle")
        if subtitle:
            sub_para = doc.add_paragraph()
            sub_run = sub_para.add_run(subtitle)
            sub_run.font.italic = True
            sub_run.font.size = DocxPt(11)
            sub_run.font.color.rgb = DocxRGB(0x55, 0x55, 0x55)

        # Key points subheading
        body_points = slide.get("body_points", [])
        if body_points:
            doc.add_heading("Key Points", level=2)
            for point in body_points:
                bp = doc.add_paragraph(style='List Bullet')
                bp_run = bp.add_run(point)
                bp_run.font.size = DocxPt(11)
                bp.paragraph_format.space_after = DocxPt(4)

        # Speaker notes as detailed discussion
        notes = slide.get("speaker_notes")
        if notes:
            doc.add_heading("Discussion", level=2)
            np = doc.add_paragraph()
            np.paragraph_format.first_line_indent = Cm(0.5)
            nr = np.add_run(notes)
            nr.font.size = DocxPt(11)
            nr.font.color.rgb = DocxRGB(0x33, 0x33, 0x33)

        # Add image with professional framing
        image_path = slide.get("image_path")
        if image_path and os.path.exists(image_path):
            try:
                doc.add_paragraph()  # spacing
                pic_para = doc.add_paragraph()
                pic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                doc.add_picture(image_path, width=DocxInches(5.0))
                # Caption
                caption = slide.get("image_caption", "")
                if caption:
                    cap_para = doc.add_paragraph()
                    cap_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    cap_para.paragraph_format.space_before = DocxPt(4)
                    cap_run = cap_para.add_run(caption[:200])
                    cap_run.font.italic = True
                    cap_run.font.size = DocxPt(9)
                    cap_run.font.color.rgb = DocxRGB(0x66, 0x66, 0x66)
            except Exception as e:
                logger.warning(f"Failed to add image to doc: {e}")

        # Thin separator between sections
        sep = doc.add_paragraph()
        sep.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sep.paragraph_format.space_before = DocxPt(12)
        sep.paragraph_format.space_after = DocxPt(6)
        sep_run = sep.add_run("—" * 30)
        sep_run.font.color.rgb = DocxRGB(0xDD, 0xDD, 0xDD)
        sep_run.font.size = DocxPt(8)

    # ====== FOOTER ON LAST PAGE ===============================================
    doc.add_paragraph()
    footer_para = doc.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = footer_para.add_run("Generated by AI Research PPT Generator")
    fr.font.size = DocxPt(9)
    fr.font.color.rgb = DocxRGB(0xAA, 0xAA, 0xAA)
    fr.font.italic = True

    # Save
    doc_path = os.path.join(settings.output_dir, f"{session_id}.docx")
    doc.save(doc_path)
    logger.info(f"[{session_id}] Word document saved to {doc_path}")
    return doc_path


def _sanitize_for_pdf(text: str) -> str:
    """Convert Unicode math/special characters to ASCII-safe equivalents for reportlab."""
    if not text:
        return text
    replacements = {
        # Math operators
        "√": "sqrt", "∛": "cbrt", "∑": "sum", "∏": "prod",
        "∫": "integral", "∂": "d", "∇": "nabla", "∆": "delta",
        "∞": "inf", "≈": "~=", "≠": "!=", "≤": "<=", "≥": ">=",
        "±": "+/-", "×": "x", "÷": "/", "·": "*",
        "∈": "in", "∉": "not in", "⊂": "subset", "⊆": "subset=",
        "∩": "intersect", "∪": "union", "∅": "{}",
        # Superscripts
        "⁰": "^0", "¹": "^1", "²": "^2", "³": "^3", "⁴": "^4",
        "⁵": "^5", "⁶": "^6", "⁷": "^7", "⁸": "^8", "⁹": "^9",
        "ⁿ": "^n", "ᵀ": "^T", "ᵃ": "^a", "ᵇ": "^b", "ᶜ": "^c",
        # Subscripts
        "₀": "_0", "₁": "_1", "₂": "_2", "₃": "_3", "₄": "_4",
        "₅": "_5", "₆": "_6", "₇": "_7", "₈": "_8", "₉": "_9",
        "ₐ": "_a", "ₑ": "_e", "ₒ": "_o", "ₙ": "_n", "ₖ": "_k",
        "ᵢ": "_i", "ⱼ": "_j", "ᵥ": "_v", "ₓ": "_x",
        # Greek letters (common in ML papers)
        "α": "alpha", "β": "beta", "γ": "gamma", "δ": "delta",
        "ε": "epsilon", "ζ": "zeta", "η": "eta", "θ": "theta",
        "λ": "lambda", "μ": "mu", "ν": "nu", "ξ": "xi",
        "π": "pi", "ρ": "rho", "σ": "sigma", "τ": "tau",
        "φ": "phi", "ψ": "psi", "ω": "omega",
        "Γ": "Gamma", "Δ": "Delta", "Θ": "Theta", "Λ": "Lambda",
        "Σ": "Sigma", "Φ": "Phi", "Ψ": "Psi", "Ω": "Omega",
        # Arrows
        "→": "->", "←": "<-", "↔": "<->", "⇒": "=>", "⇐": "<=",
        "↑": "^", "↓": "v",
        # Quotes and dashes
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u2013": "-", "\u2014": "--", "\u2026": "...",
        # Other
        "\u00b7": "*", "\u2022": "*",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    # Final pass: replace any remaining non-ASCII with '?'
    return text.encode("ascii", errors="replace").decode("ascii")


def _generate_pdf_document(session_id: str, slides: list, paper: dict) -> str:
    """Generate a professional PDF report from the slide content."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.lib.units import cm, mm
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
        PageBreak, Table, TableStyle, HRFlowable,
    )
    from reportlab.lib import colors
    from PIL import Image as PILImage

    pdf_path = os.path.join(settings.output_dir, f"{session_id}.pdf")

    # Colors
    NAVY = HexColor("#1A237E")
    ACCENT = HexColor("#1976D2")
    DARK_TEXT = HexColor("#212121")
    GRAY_TEXT = HexColor("#757575")
    LIGHT_LINE = HexColor("#E0E0E0")

    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=A4,
        topMargin=2.5 * cm,
        bottomMargin=2.5 * cm,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
    )

    width = A4[0] - 5 * cm  # available text width

    styles = getSampleStyleSheet()

    # Custom styles
    styles.add(ParagraphStyle(
        'DocTitle', parent=styles['Title'],
        fontSize=26, leading=32, textColor=NAVY,
        alignment=TA_CENTER, spaceAfter=10,
        fontName='Helvetica-Bold',
    ))
    styles.add(ParagraphStyle(
        'DocSubtitle', parent=styles['Normal'],
        fontSize=13, leading=18, textColor=GRAY_TEXT,
        alignment=TA_CENTER, spaceAfter=6,
        fontName='Helvetica-Oblique',
    ))
    styles.add(ParagraphStyle(
        'DocMeta', parent=styles['Normal'],
        fontSize=10, leading=14, textColor=GRAY_TEXT,
        alignment=TA_CENTER, spaceAfter=4,
    ))
    styles.add(ParagraphStyle(
        'SectionHead', parent=styles['Heading1'],
        fontSize=16, leading=22, textColor=NAVY,
        spaceBefore=20, spaceAfter=8,
        fontName='Helvetica-Bold',
    ))
    styles.add(ParagraphStyle(
        'SubHead', parent=styles['Heading2'],
        fontSize=12, leading=16, textColor=ACCENT,
        spaceBefore=12, spaceAfter=4,
        fontName='Helvetica-Bold',
    ))
    styles.add(ParagraphStyle(
        'Body', parent=styles['Normal'],
        fontSize=10.5, leading=15, textColor=DARK_TEXT,
        alignment=TA_JUSTIFY, spaceAfter=6,
        fontName='Helvetica',
    ))
    styles.add(ParagraphStyle(
        'BulletText', parent=styles['Normal'],
        fontSize=10.5, leading=15, textColor=DARK_TEXT,
        leftIndent=20, spaceAfter=4,
        fontName='Helvetica',
        bulletIndent=8,
    ))
    styles.add(ParagraphStyle(
        'Caption', parent=styles['Normal'],
        fontSize=9, leading=12, textColor=GRAY_TEXT,
        alignment=TA_CENTER, spaceAfter=10,
        fontName='Helvetica-Oblique',
    ))
    styles.add(ParagraphStyle(
        'Footer', parent=styles['Normal'],
        fontSize=8, textColor=GRAY_TEXT,
        alignment=TA_CENTER,
    ))

    story = []

    title_text = paper.get("title", "Research Paper Analysis")
    authors = paper.get("authors", [])
    arxiv_id = paper.get("arxiv_id", "")
    abstract = paper.get("abstract", "")

    # ====== TITLE PAGE ========================================================
    story.append(Spacer(1, 6 * cm))
    story.append(HRFlowable(width="60%", color=ACCENT, thickness=2, spaceAfter=15))
    story.append(Paragraph(_sanitize_for_pdf(title_text), styles['DocTitle']))
    story.append(Paragraph("Research Paper Analysis", styles['DocSubtitle']))
    story.append(HRFlowable(width="60%", color=ACCENT, thickness=2, spaceBefore=15, spaceAfter=20))
    if authors:
        story.append(Paragraph(_sanitize_for_pdf(', '.join(authors[:10])), styles['DocMeta']))
    meta_parts = []
    if arxiv_id:
        meta_parts.append(f"ArXiv: {arxiv_id}")
    meta_parts.append(f"Generated: {date.today().strftime('%B %d, %Y')}")
    story.append(Paragraph(" | ".join(meta_parts), styles['DocMeta']))
    story.append(PageBreak())

    # ====== TABLE OF CONTENTS =================================================
    story.append(Paragraph("Table of Contents", styles['SectionHead']))
    story.append(Spacer(1, 8))
    section_num = 0
    for slide in slides:
        if slide.get("slide_type") == "title":
            continue
        section_num += 1
        toc_text = f'<font color="#1976D2"><b>{section_num}.</b></font>  {_sanitize_for_pdf(slide.get("title", "Section"))}'
        story.append(Paragraph(toc_text, styles['Body']))
    story.append(PageBreak())

    # ====== ABSTRACT ==========================================================
    if abstract:
        story.append(Paragraph("Abstract", styles['SectionHead']))
        story.append(Paragraph(_sanitize_for_pdf(abstract), styles['Body']))
        story.append(Spacer(1, 12))

    # ====== SECTIONS ==========================================================
    section_num = 0
    for slide in slides:
        if slide.get("slide_type") == "title":
            continue

        section_num += 1
        slide_title = slide.get("title", "Section")

        story.append(Paragraph(_sanitize_for_pdf(f"{section_num}. {slide_title}"), styles['SectionHead']))

        subtitle = slide.get("subtitle")
        if subtitle:
            story.append(Paragraph(f"<i>{_sanitize_for_pdf(subtitle)}</i>", styles['DocSubtitle']))

        body_points = slide.get("body_points", [])
        if body_points:
            story.append(Paragraph("Key Points", styles['SubHead']))
            for point in body_points:
                safe_point = _sanitize_for_pdf(point).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                story.append(Paragraph(f"•  {safe_point}", styles['BulletText']))

        notes = slide.get("speaker_notes")
        if notes:
            story.append(Paragraph("Discussion", styles['SubHead']))
            safe_notes = _sanitize_for_pdf(notes).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe_notes, styles['Body']))

        # Image
        image_path = slide.get("image_path")
        if image_path and os.path.exists(image_path):
            try:
                with PILImage.open(image_path) as pil_img:
                    iw, ih = pil_img.size
                max_w = width * 0.85
                max_h = 12 * cm
                scale = min(max_w / iw, max_h / ih, 1.0)
                story.append(Spacer(1, 8))
                story.append(RLImage(image_path, width=iw * scale, height=ih * scale))
                caption = slide.get("image_caption", "")
                if caption:
                    safe_cap = _sanitize_for_pdf(caption[:200]).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    story.append(Paragraph(safe_cap, styles['Caption']))
            except Exception as e:
                logger.warning(f"Failed to add image to PDF: {e}")

        # Section separator
        story.append(Spacer(1, 8))
        story.append(HRFlowable(width="40%", color=LIGHT_LINE, thickness=0.5, spaceAfter=6))

    # Footer
    story.append(Spacer(1, 20))
    story.append(Paragraph("Generated by AI Research PPT Generator", styles['Footer']))

    # Build PDF
    doc.build(story)
    logger.info(f"[{session_id}] PDF document saved to {pdf_path}")
    return pdf_path


def ppt_generation_node(state: ResearchState) -> dict:
    """Generate comprehensive PPT, Word document, and PDF from approved slides."""
    session_id = state["session_id"]
    approved_slides = state.get("approved_slides") or state.get("slide_contents", [])
    logger.info(f"[{session_id}] Generating PPT with {len(approved_slides)} slides")

    prs = create_presentation()
    total = len(approved_slides)

    for i, content in enumerate(approved_slides, start=1):
        slide_type = content.get("slide_type", "content")

        try:
            handler = SLIDE_HANDLERS.get(slide_type, _add_content_slide_with_image)
            handler(prs, content, i, total)
        except Exception as e:
            logger.error(f"[{session_id}] Error adding slide {i} ({slide_type}): {e}")
            _add_content_slide_with_image(prs, content, i, total)

    # Save PPT
    os.makedirs(settings.output_dir, exist_ok=True)
    file_path = os.path.join(settings.output_dir, f"{session_id}.pptx")
    prs.save(file_path)

    # Generate companion Word document
    paper = state.get("processed_paper") or state.get("selected_paper") or {}
    doc_path = None
    try:
        doc_path = _generate_word_document(session_id, approved_slides, paper)
    except Exception as e:
        logger.warning(f"[{session_id}] Word document generation failed: {e}")

    # Generate companion PDF document
    pdf_path = None
    try:
        pdf_path = _generate_pdf_document(session_id, approved_slides, paper)
    except Exception as e:
        logger.warning(f"[{session_id}] PDF document generation failed: {e}")

    title = paper.get("title", "AI Research Presentation")

    generated_ppt: GeneratedPPT = {
        "file_path": file_path,
        "doc_path": doc_path,
        "pdf_path": pdf_path,
        "session_id": session_id,
        "slide_count": len(approved_slides),
        "topics_covered": [title[:100]],
        "generated_at": date.today().isoformat(),
    }

    logger.info(f"[{session_id}] PPT saved to {file_path} with {len(approved_slides)} slides")

    # Clean up extracted assets (PDFs, images) - they're now embedded in pptx/docx/pdf
    assets_dir = os.path.join(settings.output_dir, session_id, "assets")
    if os.path.isdir(assets_dir):
        import shutil
        shutil.rmtree(assets_dir, ignore_errors=True)
        logger.info(f"[{session_id}] Cleaned up assets directory")

    return {
        "current_stage": Stage.AWAITING_FINAL_REVIEW,
        "generated_ppt": generated_ppt,
    }
