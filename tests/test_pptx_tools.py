"""Tests for PPTX generation tools."""
import os
import tempfile
import pytest
from pptx.util import Inches
from pptx.dml.color import RGBColor

from backend.tools.pptx_tools import (
    create_presentation,
    set_slide_background,
    add_header_bar,
    add_bullet_points,
    add_page_number,
    add_footer,
    add_accent_line,
    DARK_BLUE,
    LIGHT_GRAY,
)


class TestCreatePresentation:
    def test_creates_widescreen(self):
        prs = create_presentation()
        assert prs is not None
        assert prs.slide_width == Inches(13.33)
        assert prs.slide_height == Inches(7.5)

    def test_save_and_load(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            prs.save(f.name)
            assert os.path.exists(f.name)
            assert os.path.getsize(f.name) > 1000
            os.unlink(f.name)


class TestSlideBackground:
    def test_set_background(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide, DARK_BLUE)
        # Should not raise


class TestHeaderBar:
    def test_add_header(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_bar(slide, "Test Title", "Test Subtitle")
        assert len(slide.shapes) >= 2  # header rect + title text

    def test_add_header_no_subtitle(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_header_bar(slide, "Title Only")
        assert len(slide.shapes) >= 1


class TestBulletPoints:
    def test_add_bullets(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        points = ["First point about the topic", "Second important detail", "Third conclusion"]
        add_bullet_points(slide, points)
        assert len(slide.shapes) >= 1

    def test_empty_bullets(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bullet_points(slide, [])
        # Should not crash


class TestPageNumber:
    def test_add_page_number(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_page_number(slide, 1, 10)
        assert len(slide.shapes) >= 1


class TestFooter:
    def test_add_footer(self):
        prs = create_presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_footer(slide, "Test Footer Text")
        assert len(slide.shapes) >= 1
